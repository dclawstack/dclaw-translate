from __future__ import annotations

import uuid
from pathlib import Path

from fastapi import UploadFile
from sqlalchemy.ext.asyncio import AsyncSession

from app.models.document import Document
from app.repositories.document_repo import DocumentRepository
from app.services.llm_service import LLMService

SUPPORTED_TYPES = {"pdf", "docx", "pptx"}


class DocumentService:
    def __init__(self, settings) -> None:
        self.upload_dir = Path(settings.upload_dir)
        self.upload_dir.mkdir(parents=True, exist_ok=True)

    async def save_upload(self, file: UploadFile) -> tuple[str, str, int]:
        """Save uploaded file. Returns (saved_path, file_type, file_size)."""
        original_name = file.filename or "upload"
        ext = Path(original_name).suffix.lstrip(".").lower()
        if ext not in SUPPORTED_TYPES:
            raise ValueError(f"Unsupported file type: .{ext}. Supported: pdf, docx, pptx")
        unique_name = f"{uuid.uuid4()}_{original_name}"
        dest = self.upload_dir / unique_name
        content = await file.read()
        dest.write_bytes(content)
        return str(dest), ext, len(content)

    def extract_text(self, path: str, file_type: str) -> tuple[str, int]:
        """Extract text from document. Returns (text, page_count)."""
        try:
            if file_type == "docx":
                from docx import Document as DocxDocument  # type: ignore
                doc = DocxDocument(path)
                text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
                page_count = max(1, len(doc.paragraphs) // 20)
                return text, page_count
            elif file_type == "pdf":
                import fitz  # type: ignore
                doc = fitz.open(path)
                pages = []
                for page in doc:
                    pages.append(page.get_text())
                doc.close()
                return "\n".join(pages), len(pages)
            elif file_type == "pptx":
                from pptx import Presentation  # type: ignore
                prs = Presentation(path)
                slides_text = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            slides_text.append(shape.text)
                return "\n".join(slides_text), len(prs.slides)
        except Exception:
            return "", 0
        return "", 0

    async def translate_document(
        self,
        db: AsyncSession,
        doc_id: uuid.UUID,
        llm_service: LLMService,
    ) -> Document:
        repo = DocumentRepository(db)
        doc = await repo.get_by_id(doc_id)
        if doc is None:
            raise ValueError(f"Document {doc_id} not found")

        try:
            doc.status = "extracting"
            await repo.update(doc)

            text, page_count = self.extract_text(doc.original_path, doc.file_type)
            doc.page_count = page_count
            doc.word_count = len(text.split()) if text else 0
            doc.status = "translating"
            await repo.update(doc)

            # Chunk text into ~1000-char pieces and translate each
            system_prompt = (
                f"You are a professional translator. Translate from {doc.source_lang} "
                f"to {doc.target_lang}. Return ONLY the translated text, nothing else."
            )
            if text:
                chunks = [text[i:i + 1000] for i in range(0, len(text), 1000)]
                translated_chunks = []
                for chunk in chunks:
                    result = await llm_service.translate(
                        db,
                        chunk,
                        doc.source_lang,
                        doc.target_lang,
                        system_prompt,
                    )
                    translated_chunks.append(result["translated_text"])
                translated_text = "\n".join(translated_chunks)
            else:
                translated_text = ""

            doc.status = "reconstructing"
            await repo.update(doc)

            # Save translated text as .translated.txt alongside original
            translated_path = doc.original_path + ".translated.txt"
            Path(translated_path).write_text(translated_text, encoding="utf-8")
            doc.translated_path = translated_path
            doc.status = "completed"
            await repo.update(doc)

        except Exception as e:
            doc.status = "failed"
            doc.error_message = str(e)
            await repo.update(doc)

        return doc

    async def delete_files(self, doc: Document) -> None:
        """Delete original and translated files from filesystem."""
        for path_str in [doc.original_path, doc.translated_path]:
            if path_str:
                p = Path(path_str)
                if p.exists():
                    p.unlink()
